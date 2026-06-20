"""
Phase 1 AI pipeline. Orchestrates the per-item processing:

  raw_text / fetched_url_body
        │
        ▼
  classify_item (gpt-4o)   →   title, summary, tags, entities, confidence, folder
  summarise (gpt-4o-mini)  →   richer summary if classification one was sparse
  embed_item               →   1536-dim vector
        │
        ▼
  UPDATE items SET ...
        │
        ▼
  sync_item_to_neo4j.delay   (existing task — generates edges)

GitHub Models rate limits apply (10–15 RPM on free tier), so we run sequentially
within a single task rather than fanning out parallel calls.
"""
from __future__ import annotations

import asyncio
import logging
import uuid
from datetime import datetime

import httpx
from bs4 import BeautifulSoup
from celery import shared_task
from sqlalchemy import update, select

from app.tasks.celery_app import celery_app
from app.tasks._db import task_session
from app.models.item import Item
from app.models.folder import Folder
from app.ai.llm import classify_item, summarise
from app.ai.embeddings import embed_item

log = logging.getLogger(__name__)

MAX_TEXT_FOR_AI = 8000   # chars sent to classifier
URL_TIMEOUT = 15
MAX_FOLDER_DEPTH = 2     # 0-indexed; 3 levels total (matches folders router cap)


# ── Content extraction ────────────────────────────────────────────────────────

async def _fetch_url_text(url: str) -> str:
    """Fetch a URL and return its main body text (rough — strips tags only)."""
    try:
        async with httpx.AsyncClient(timeout=URL_TIMEOUT, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 SecondBrain/1.0"})
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "lxml")
            for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
                tag.decompose()
            text = soup.get_text(" ", strip=True)
            return text[:MAX_TEXT_FOR_AI]
    except Exception as exc:
        log.warning("URL fetch failed for %s: %s", url, exc)
        return ""


def _extract_file_text(item: Item) -> str:
    """Pull text out of an uploaded file for the classifier.
    PDF/DOCX/XLSX/PPTX go through their respective libraries; text files are
    decoded directly; images/audio/video yield empty string (handled separately)."""
    from app import storage

    if not item.storage_key or not storage.file_exists(item.storage_key):
        return ""
    try:
        raw = storage.read_file(item.storage_key)
    except Exception as exc:
        log.warning("Could not read file for %s: %s", item.id, exc)
        return ""

    if item.content_type == "pdf":
        try:
            import io
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            pages = [page.extract_text() or "" for page in reader.pages]
            return "\n".join(pages)[:MAX_TEXT_FOR_AI]
        except Exception as exc:
            log.warning("PDF extraction failed for %s: %s", item.id, exc)
            return ""

    if item.content_type == "doc":
        filename = (item.metadata_ or {}).get("filename", "")
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext in ("docx",):
            try:
                import io
                from docx import Document
                doc = Document(io.BytesIO(raw))
                text = "\n".join(p.text for p in doc.paragraphs)
                return text[:MAX_TEXT_FOR_AI]
            except Exception as exc:
                log.warning("DOCX extraction failed for %s: %s", item.id, exc)

        if ext in ("xlsx", "xls"):
            try:
                import io, openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
                lines = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        cell_vals = [str(c) for c in row if c is not None]
                        if cell_vals:
                            lines.append("\t".join(cell_vals))
                return "\n".join(lines)[:MAX_TEXT_FOR_AI]
            except Exception as exc:
                log.warning("XLSX extraction failed for %s: %s", item.id, exc)

        if ext in ("pptx",):
            try:
                import io
                from pptx import Presentation
                prs = Presentation(io.BytesIO(raw))
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            lines.append(shape.text)
                return "\n".join(lines)[:MAX_TEXT_FOR_AI]
            except Exception as exc:
                log.warning("PPTX extraction failed for %s: %s", item.id, exc)

        # fallback: try raw UTF-8 decode (CSV, plain text uploaded as doc)
        try:
            return raw.decode("utf-8", errors="ignore")[:MAX_TEXT_FOR_AI]
        except Exception:
            return ""

    if item.content_type in ("text",):
        try:
            return raw.decode("utf-8", errors="ignore")[:MAX_TEXT_FOR_AI]
        except Exception:
            return ""

    # image / audio / video — no text layer extracted here
    return ""


async def _describe_image(item: Item) -> str:
    """Send image to GitHub Models vision to get a text description."""
    from app import storage
    import base64

    if not item.storage_key or not storage.file_exists(item.storage_key):
        return ""
    try:
        raw = storage.read_file(item.storage_key)
    except Exception as exc:
        log.warning("Cannot read image for %s: %s", item.id, exc)
        return ""

    filename = (item.metadata_ or {}).get("filename", "image.jpg")
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "jpg"
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "webp": "image/webp", "gif": "image/gif"}
    mime = mime_map.get(ext, "image/jpeg")

    if len(raw) > 5 * 1024 * 1024:
        log.info("Image %s is >5MB, skipping vision description", item.id)
        return ""

    b64 = base64.b64encode(raw).decode()
    data_url = f"data:{mime};base64,{b64}"

    from app.config import settings
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            res = await client.post(
                f"{settings.github_models_endpoint}/chat/completions",
                headers={
                    "Authorization": f"Bearer {settings.github_token}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": "openai/gpt-4o",
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": (
                                        "Describe this image in detail. Include: what you see, "
                                        "any text visible, objects, scenes, colours, people, places. "
                                        "Also suggest 3-6 relevant tags. Be concise (under 200 words)."
                                    ),
                                },
                                {"type": "image_url", "image_url": {"url": data_url}},
                            ],
                        }
                    ],
                    "max_tokens": 400,
                },
            )
            res.raise_for_status()
            return res.json()["choices"][0]["message"]["content"]
    except Exception as exc:
        log.warning("Vision description failed for %s: %s", item.id, exc)
        return ""


def _pick_content(item: Item, fetched: str, image_desc: str = "") -> str:
    """Choose the best text representation of an item for AI processing."""
    if item.raw_text:
        return item.raw_text[:MAX_TEXT_FOR_AI]
    if fetched:
        return fetched
    if image_desc:
        return image_desc
    extracted = _extract_file_text(item)
    if extracted.strip():
        return extracted
    if item.title:
        return item.title
    return ""


# ── AI folder routing ───────────────────────────────────────────────────────

async def _load_folders(db, user_id: str) -> list[Folder]:
    return list(
        (await db.execute(select(Folder).where(Folder.user_id == user_id))).scalars().all()
    )


def _folder_tree_json(folders: list[Folder]) -> str:
    """Compact JSON of the user's folders so the classifier can reuse existing
    ones instead of always inventing new paths."""
    import json as _json
    by_id = {f.id: f for f in folders}

    def path_of(f: Folder) -> list[str]:
        chain, cur, guard = [], f, 0
        while cur is not None and guard < 5:
            chain.append(cur.name)
            cur = by_id.get(cur.parent_id) if cur.parent_id else None
            guard += 1
        return list(reversed(chain))

    return _json.dumps([path_of(f) for f in folders], ensure_ascii=False)


async def _resolve_folder_path(
    db, user_id: str, path: list[str]
) -> "uuid.UUID | None":
    """Map an AI-suggested folder path (e.g. ["Learning", "Group Theory"]) to a
    real folder id. Reuses existing folders (case-insensitive) at each level and
    creates the missing tail segments (marked ai_generated). Returns the leaf id.

    Strategy: prefer existing, create if needed (capped at 3 levels)."""
    import uuid as _uuid

    clean = [seg.strip() for seg in (path or []) if seg and seg.strip()]
    if not clean:
        return None
    clean = clean[: MAX_FOLDER_DEPTH + 1]  # never exceed the 3-level cap

    folders = await _load_folders(db, user_id)
    # index existing folders by (parent_id, lowercased name) for O(1) level lookup
    index: dict[tuple, Folder] = {
        (f.parent_id, f.name.strip().lower()): f for f in folders
    }

    parent_id = None
    leaf_id = None
    for depth, segment in enumerate(clean):
        key = (parent_id, segment.lower())
        existing = index.get(key)
        if existing is not None:
            parent_id = existing.id
            leaf_id = existing.id
            continue
        # create this segment
        new_folder = Folder(
            user_id=_uuid.UUID(user_id),
            parent_id=parent_id,
            name=segment,
            depth=depth,
            ai_generated=True,
        )
        db.add(new_folder)
        await db.flush()  # assign id
        index[key] = new_folder
        parent_id = new_folder.id
        leaf_id = new_folder.id
        log.info("AI created folder '%s' (depth=%d) for user %s", segment, depth, user_id)

    return leaf_id


# ── Main task ─────────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True,
    max_retries=3,
    default_retry_delay=60,
    name="app.tasks.ai_processing.process_item",
)
def process_item(self, item_id: str, user_id: str) -> None:
    """Full AI pipeline for one item, then trigger Neo4j sync."""
    try:
        asyncio.run(_process(item_id, user_id))
    except Exception as exc:
        log.error("process_item failed for %s: %s", item_id, exc, exc_info=True)
        raise self.retry(exc=exc, countdown=60 * (2 ** self.request.retries))


async def _process(item_id: str, user_id: str) -> None:
    async with task_session() as db:
        item = (
            await db.execute(select(Item).where(Item.id == item_id))
        ).scalar_one_or_none()
        if not item:
            log.warning("process_item: %s not found", item_id)
            return

        fetched = ""
        if item.content_type == "url" and item.source_url:
            fetched = await _fetch_url_text(item.source_url)

        image_desc = ""
        if item.content_type == "image":
            image_desc = await _describe_image(item)

        content = _pick_content(item, fetched, image_desc)
        if not content.strip():
            log.info("process_item: %s has no content, skipping AI", item_id)
            await _trigger_sync(item_id, user_id)
            return

        meta = item.metadata_ or {}
        user_set_tags = bool(meta.get("user_set_tags"))
        ai_organise = bool(meta.get("ai_organise"))

        # Only feed the existing folder tree to the classifier when we actually
        # intend to auto-file — saves tokens otherwise.
        folder_tree_json = "[]"
        if ai_organise and item.folder_id is None:
            folders = await _load_folders(db, str(item.user_id))
            folder_tree_json = _folder_tree_json(folders)

        classification = await classify_item(
            content_type=item.content_type,
            content_text=content,
            folder_tree_json=folder_tree_json,
            skip_tags=user_set_tags,
        )

        summary = classification.get("summary") or ""
        if not summary:
            summary = await summarise(content)

        final_tags = item.tags if user_set_tags else (classification.get("tags") or [])

        embedding = await embed_item({
            "title": classification.get("title") or item.title,
            "summary": summary,
            "raw_text": content,
            "tags": final_tags,
        })

        confidence = float(classification.get("confidence") or 0.0)

        # AI folder routing — only when the user opted in AND left the folder
        # blank (an explicit manual choice always wins).
        resolved_folder_id = None
        if ai_organise and item.folder_id is None:
            resolved_folder_id = await _resolve_folder_path(
                db, str(item.user_id), classification.get("suggested_folder") or []
            )

        values = dict(
            ai_title=classification.get("title") or None,
            summary=summary or None,
            tags=final_tags,
            entities=classification.get("entities") or {},
            confidence=confidence,
            needs_review=confidence < 0.5,
            embedding=embedding,
            indexed_at=datetime.utcnow(),
            raw_text=content if not item.raw_text else item.raw_text,
        )
        if resolved_folder_id is not None:
            values["folder_id"] = resolved_folder_id

        await db.execute(update(Item).where(Item.id == item_id).values(**values))
        await db.commit()
        log.info(
            "AI processing complete for %s (conf=%.2f, folder=%s)",
            item_id, confidence, resolved_folder_id,
        )

        # Sync the resolved folder name to the Neo4j node (used by cluster mode).
        if resolved_folder_id is not None:
            folder = await db.get(Folder, resolved_folder_id)
            if folder is not None:
                from app.tasks.sync import update_item_folder
                update_item_folder.delay(
                    str(item_id), folder.name, str(resolved_folder_id)
                )

    await _trigger_sync(item_id, user_id)


async def _trigger_sync(item_id: str, user_id: str) -> None:
    """Hand off to the existing Neo4j sync + edge generation task."""
    from app.tasks.sync import sync_item_to_neo4j
    sync_item_to_neo4j.delay(item_id, user_id)
